"""Utilities for rendering data frames as Great Tables inside PowerPoint decks."""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Mapping, Optional, Sequence, Union
import tempfile

try:
    import pandas as pd
except ModuleNotFoundError as exc:  # pragma: no cover - pandas is required at runtime.
    raise ModuleNotFoundError(
        "pandas is required to use powerpoint-lib. Install it with `pip install pandas`."
    ) from exc

try:
    from pptx import Presentation
    from pptx.util import Inches
except ModuleNotFoundError as exc:  # pragma: no cover - python-pptx is required at runtime.
    raise ModuleNotFoundError(
        "python-pptx is required to use powerpoint-lib. Install it with `pip install python-pptx`."
    ) from exc

try:
    from great_tables import GT, loc, style
except ModuleNotFoundError as exc:  # pragma: no cover - great-tables is required at runtime.
    raise ModuleNotFoundError(
        "great-tables is required to use powerpoint-lib. Install it with `pip install great-tables`."
    ) from exc


HexColor = str


@dataclass(slots=True)
class ExportOptions:
    """Configuration for rendering a DataFrame as a PowerPoint slide."""

    font: str = "Source Sans Pro"
    header_fill: HexColor = "#2F5597"
    header_font_color: HexColor = "#FFFFFF"
    row_fill: HexColor = "#FFFFFF"
    row_font_color: HexColor = "#222222"
    alternating_fill: Optional[HexColor] = "#F2F2F2"
    top_inches: float = 1.0
    left_inches: float = 0.75
    width_inches: Optional[float] = 9.0


class PowerPointTableExporter:
    """Render pandas data frames as table images and embed them into PowerPoint decks."""

    def __init__(self, options: Optional[ExportOptions] = None) -> None:
        self.options = options or ExportOptions()

    def export(
        self,
        dataframe: Union[pd.DataFrame, "pd.io.formats.style.Styler"],
        template_path: Union[str, Path],
        output_path: Union[str, Path],
        *,
        formatters: Optional[Mapping[str, Union[str, Callable[[object], str]]]] = None,
        column_order: Optional[Sequence[str]] = None,
    ) -> Path:
        """Export a DataFrame to a PowerPoint deck.

        Parameters
        ----------
        dataframe:
            The pandas DataFrame (or Styler) to render as a table.
        template_path:
            Path to an existing PowerPoint template to use as the base presentation.
        output_path:
            The final path for the rendered presentation.
        formatters:
            Optional column-to-format mapping passed to :meth:`pandas.DataFrame.style.format`
            to respect custom display strings when a plain DataFrame (not a Styler) is
            supplied. Each value can be a format string (e.g. ``"${:.0f}"``) or a callable
            accepting a cell value and returning a string.
        column_order:
            Optional sequence defining the order of columns to include in the table.
        """

        template = Path(template_path).expanduser().resolve()
        if not template.exists():
            raise FileNotFoundError(f"PowerPoint template not found: {template}")

        destination = Path(output_path).expanduser().resolve()
        if not destination.parent.exists():
            destination.parent.mkdir(parents=True, exist_ok=True)

        display_df = _create_display_dataframe(dataframe, formatters=formatters)
        if column_order is not None:
            missing = [column for column in column_order if column not in display_df.columns]
            if missing:
                raise KeyError(f"Columns not present in DataFrame: {missing}")
            display_df = display_df[list(column_order)]

        image_path = _render_table_image(display_df, self.options)
        try:
            self._embed_image(template, destination, image_path)
        finally:
            image_path.unlink(missing_ok=True)

        return destination

    def _embed_image(self, template: Path, destination: Path, image_path: Path) -> None:
        presentation = Presentation(str(template))
        if not presentation.slide_layouts:
            raise ValueError("The template presentation does not define any slide layouts.")

        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        options = self.options
        left = Inches(options.left_inches)
        top = Inches(options.top_inches)
        if options.width_inches is not None:
            slide.shapes.add_picture(str(image_path), left, top, width=Inches(options.width_inches))
        else:
            slide.shapes.add_picture(str(image_path), left, top)
        presentation.save(str(destination))


def export_table_to_pptx(
    dataframe: Union[pd.DataFrame, "pd.io.formats.style.Styler"],
    template_path: Union[str, Path],
    output_path: Union[str, Path],
    *,
    font: str = "Source Sans Pro",
    header_fill: HexColor = "#2F5597",
    header_font_color: HexColor = "#FFFFFF",
    row_fill: HexColor = "#FFFFFF",
    row_font_color: HexColor = "#222222",
    alternating_fill: Optional[HexColor] = "#F2F2F2",
    formatters: Optional[Mapping[str, Union[str, Callable[[object], str]]]] = None,
    column_order: Optional[Sequence[str]] = None,
) -> Path:
    """Convenience wrapper around :class:`PowerPointTableExporter`.

    Returns
    -------
    Path
        The absolute path to the rendered PowerPoint deck.
    """

    exporter = PowerPointTableExporter(
        ExportOptions(
            font=font,
            header_fill=header_fill,
            header_font_color=header_font_color,
            row_fill=row_fill,
            row_font_color=row_font_color,
            alternating_fill=alternating_fill,
        )
    )
    return exporter.export(
        dataframe=dataframe,
        template_path=template_path,
        output_path=output_path,
        formatters=formatters,
        column_order=column_order,
    )


def _create_display_dataframe(
    dataframe: Union[pd.DataFrame, "pd.io.formats.style.Styler"],
    *,
    formatters: Optional[Mapping[str, Union[str, Callable[[object], str]]]],
) -> pd.DataFrame:
    """Return a DataFrame containing the formatted display strings."""

    if _is_styler(dataframe):
        return _extract_display_values_from_styler(dataframe)  # type: ignore[arg-type]

    df = dataframe.copy()
    if formatters:
        styler = df.style.format(formatters)
        return _extract_display_values_from_styler(styler)

    object_df = df.copy()
    for column in object_df.columns:
        object_df[column] = object_df[column].apply(_format_cell_value)
    return object_df


def _is_styler(obj: object) -> bool:
    return obj.__class__.__name__ == "Styler"


def _extract_display_values_from_styler(styler: "pd.io.formats.style.Styler") -> pd.DataFrame:
    translation = styler._translate()
    header_rows = translation.get("header", [])
    if not header_rows:
        columns = list(styler.data.columns)
    else:
        header_values = [cell.get("display_value", "") for cell in header_rows[-1]]
        if len(header_values) > len(styler.data.columns):
            header_values = header_values[-len(styler.data.columns) :]
        columns = header_values
    body_rows = []
    for row in translation.get("body", []):
        values = [cell.get("display_value", "") for cell in row]
        if len(values) > len(columns):
            values = values[-len(columns) :]
        body_rows.append(values)
    return pd.DataFrame(body_rows, columns=columns)


def _format_cell_value(value: object) -> str:
    if value is None:
        return ""
    return str(value)


def _render_table_image(dataframe: pd.DataFrame, options: ExportOptions) -> Path:
    table = GT(dataframe)
    table = table.opt_table_font(font=options.font)
    table = table.tab_style(style.fill(color=options.header_fill), loc.column_labels())
    table = table.tab_style(style.text(color=options.header_font_color), loc.column_labels())
    table = table.tab_style(style.fill(color=options.row_fill), loc.body())
    table = table.tab_style(style.text(color=options.row_font_color), loc.body())

    if options.alternating_fill:
        body_rows_selector = getattr(loc, "body_rows", None)
        body_selector = getattr(loc, "body", None)
        try:
            if callable(body_rows_selector):
                table = table.tab_style(
                    style.fill(color=options.alternating_fill),
                    body_rows_selector(groups="odd"),
                )
            elif callable(body_selector):
                table = table.tab_style(
                    style.fill(color=options.alternating_fill),
                    body_selector(rows=slice(1, None, 2)),
                )
        except Exception:  # pragma: no cover - rely on library specifics at runtime.
            pass

    with tempfile.NamedTemporaryFile(prefix="powerpoint-lib-", suffix=".png", delete=False) as handle:
        image_path = Path(handle.name)

    save_method = getattr(table, "save", None)
    export_method = getattr(table, "export", None)
    if callable(save_method):
        save_method(str(image_path))
    elif callable(export_method):
        export_method(path=str(image_path))
    else:  # pragma: no cover - depends on the external library API.
        raise AttributeError("Unable to locate a save/export method on the GT table instance.")
    return image_path
